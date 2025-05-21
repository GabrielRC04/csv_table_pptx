from pptx import Presentation
from pptx.util import Inches
import csv

# Inicialize os vetores
datas = []
horarios = []
dias = []
celebracoes = []
operadores = []

#Configuração da formação dos slides
prs = Presentation()
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes


shapes.title.text = 'Escala Equipe de Som \n Matriz São Pedro'
linhas = len(datas) + 1  # +1 para o cabeçalho
colunas = 5
left = Inches(0.2) 
top = Inches(2.0)
width = Inches(9.0)
height = Inches(2.0)

# Abra e leia o arquivo CSV
with open('escala.csv', mode='r', encoding='utf-8') as arquivo:
    leitor = csv.DictReader(arquivo)
    for linha in leitor:
        datas.append(linha['Data'])
        horarios.append(linha['Horário'])
        dias.append(linha['Dia'])
        celebracoes.append(linha['Celebração'])
        operadores.append(linha['Operador'])
        linhas+=1

table = shapes.add_table(linhas, colunas, left, top, width, height).table

# set column widths
table.columns[0].width = Inches(1.0)
table.columns[1].width = Inches(1.2)
table.columns[2].width = Inches(2.5)
table.columns[3].width = Inches(2.5)
table.columns[4].width = Inches(2.5)

# write column headings
table.cell(0, 0).text = 'Data'
table.cell(0, 1).text = 'Horário'
table.cell(0, 2).text = 'Dia'
table.cell(0, 3).text = 'Celebração'
table.cell(0, 4).text = 'Operador'


i=0
# write body cells
for i in range(len(datas)):
    table.cell(i+1, 0).text = datas[i]
    table.cell(i+1, 1).text = horarios[i]
    table.cell(i+1, 2).text = dias[i]
    table.cell(i+1, 3).text = celebracoes[i]
    table.cell(i+1, 4).text = operadores[i]

prs.save('Escala.pptx')